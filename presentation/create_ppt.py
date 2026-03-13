"""
JalSakhi - Round 2 Presentation Generator
World Water Day Ideation Challenge 2026
L&T Construction Water & Effluent Treatment
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys

# ── CONFIG ────────────────────────────────────────────────
TEAM_NAME = sys.argv[1] if len(sys.argv) > 1 else "JalSakhi"
INSTITUTE_NAME = sys.argv[2] if len(sys.argv) > 2 else "Institute Name"

# Colors
DEEP_BLUE = RGBColor(0x0D, 0x2B, 0x45)    # Primary dark
WATER_BLUE = RGBColor(0x1B, 0x6B, 0x93)    # Water themed
TEAL = RGBColor(0x00, 0xA8, 0xB5)          # Accent
LIGHT_TEAL = RGBColor(0xE0, 0xF7, 0xFA)    # Light bg
ORANGE = RGBColor(0xFF, 0x6B, 0x35)         # Highlight
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_HEADER_BG = RGBColor(0x0D, 0x2B, 0x45)
TABLE_ALT_BG = RGBColor(0xF0, 0xF8, 0xFF)
RED_ACCENT = RGBColor(0xE8, 0x3E, 0x3E)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_bg_rect(slide, color):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape


def add_accent_bar(slide, y=Inches(0), width=SLIDE_WIDTH, height=Inches(0.06), color=TEAL):
    """Add a thin accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), y, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_side_bar(slide, width=Inches(0.4), color=WATER_BLUE):
    """Add left side accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), width, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    """Add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bullet_color=TEAL, spacing=1.3):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Use bullet character
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.line_spacing = Pt(font_size * spacing)
        p.level = 0

        # Add bullet via XML
        from pptx.oxml.ns import qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u25CF'})
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': '{:02X}{:02X}{:02X}'.format(bullet_color[0], bullet_color[1], bullet_color[2]) if isinstance(bullet_color, tuple) else f'{bullet_color.red:02X}{bullet_color.green:02X}{bullet_color.blue:02X}' if hasattr(bullet_color, 'red') else '00A8B5'})
        buClr.append(srgbClr)
        # Remove existing bullet elements
        for child in list(pPr):
            if child.tag.endswith('buChar') or child.tag.endswith('buNone') or child.tag.endswith('buClr'):
                pPr.remove(child)
        pPr.append(buClr)
        pPr.append(buChar)

    return txBox


def add_stat_box(slide, left, top, width, height, number, label,
                 num_color=ORANGE, label_color=DARK_TEXT, bg_color=None):
    """Add a statistic highlight box."""
    if bg_color:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()

    add_text_box(slide, left, top + Inches(0.15), width, Inches(0.6),
                 number, font_size=32, bold=True, color=num_color,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.7), width, Inches(0.5),
                 label, font_size=12, color=label_color,
                 align=PP_ALIGN.CENTER)
    return


def add_slide_number(slide, num, total):
    """Add slide number in bottom right."""
    add_text_box(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=GRAY_TEXT,
                 align=PP_ALIGN.RIGHT)


def add_footer_bar(slide):
    """Add footer with JalSakhi branding."""
    add_accent_bar(slide, y=Inches(7.2), height=Inches(0.3), color=DEEP_BLUE)
    add_text_box(slide, Inches(0.5), Inches(7.2), Inches(5), Inches(0.3),
                 "JalSakhi | World Water Day 2026", font_size=9,
                 color=WHITE, bold=False)


def make_table(slide, left, top, width, rows_data, col_widths=None, font_size=13):
    """Create a styled table."""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0])

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top,
                                          width, Inches(0.4 * num_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_TEXT

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


# ── BUILD PRESENTATION ────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank_layout = prs.slide_layouts[6]  # Blank layout

TOTAL_SLIDES = 12


# ══════════════════════════════════════════════════════════
# TITLE SLIDE (not counted)
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg_rect(slide, DEEP_BLUE)

# Top accent line
add_accent_bar(slide, y=Inches(0), height=Inches(0.08), color=TEAL)

# Main title
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "JalSakhi", font_size=60, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(2.7), Inches(10), Inches(0.8),
             "Smartphone Electrochemical Water Forensics &\nIntelligent Treatment Advisory for Women-Led Rural Communities",
             font_size=22, color=TEAL, align=PP_ALIGN.CENTER)

# Divider line
add_accent_bar(slide, y=Inches(3.8), width=Inches(4), height=Inches(0.04), color=TEAL)
# Center the divider
shapes = slide.shapes
divider = shapes[-1]
divider.left = int((SLIDE_WIDTH - Inches(4)) / 2)

# Event info
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             "World Water Day Ideation Challenge 2026",
             font_size=20, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.4),
             'L&T Construction Water & Effluent Treatment | Theme: "Water and Gender"',
             font_size=14, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# Team info
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
             f"Team: {TEAM_NAME}", font_size=18, color=WHITE,
             align=PP_ALIGN.CENTER, bold=True)
add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.4),
             INSTITUTE_NAME, font_size=14, color=GRAY_TEXT,
             align=PP_ALIGN.CENTER)

# Bottom accent
add_accent_bar(slide, y=Inches(7.42), height=Inches(0.08), color=TEAL)


# ══════════════════════════════════════════════════════════
# SLIDE 1: PROBLEM STATEMENT - THE HOOK
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_side_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 1, TOTAL_SLIDES)

# Section label
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             "PROBLEM STATEMENT", font_size=12, bold=True, color=TEAL)

# Title
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "India Cannot Measure Its Water", font_size=36, bold=True,
             color=DEEP_BLUE)

# Story hook
hook_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.3)
)
hook_box.fill.solid()
hook_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)  # Warm orange bg
hook_box.line.fill.background()

add_text_box(slide, Inches(1.1), Inches(1.7), Inches(11), Inches(1.1),
             '"A woman in rural Bihar notices something off about her borewell water.\n'
             'The nearest accredited lab is 53 km away, costs INR 500, and takes 10 days.\n'
             'Her family drinks the water while they wait."',
             font_size=16, color=RGBColor(0x8B, 0x45, 0x00),
             align=PP_ALIGN.LEFT)

# Stats in boxes
stat_y = Inches(3.3)
stat_h = Inches(1.3)
gap = Inches(0.3)
box_w = Inches(2.5)
start_x = Inches(0.8)

stats = [
    ("1.9 M", "Rural Habitations"),
    ("2,200", "Testing Labs"),
    ("3-14 Days", "Turnaround Time"),
    ("INR 500-2000", "Cost Per Test"),
]

for i, (num, label) in enumerate(stats):
    x = start_x + i * (box_w + gap)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, stat_y, box_w, stat_h
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_TEAL
    box.line.color.rgb = TEAL
    box.line.width = Pt(1)

    add_text_box(slide, x, stat_y + Inches(0.15), box_w, Inches(0.55),
                 num, font_size=28, bold=True, color=WATER_BLUE,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, x, stat_y + Inches(0.7), box_w, Inches(0.4),
                 label, font_size=13, color=DARK_TEXT,
                 align=PP_ALIGN.CENTER)

# Gender dimension
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
             "The Gender Dimension", font_size=22, bold=True, color=DEEP_BLUE)

add_bullet_list(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5), [
    "Women & girls spend 1.4 billion hours/year collecting water in India (UNICEF)",
    "They cannot verify whether their water sources are safe",
    "By the time lab results arrive, contamination events are over - families have already consumed the water",
], font_size=15)

# Punchline
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
             '"You can\'t manage what you can\'t measure. And India cannot measure its water."',
             font_size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 2: PROPOSED SOLUTION - OVERVIEW
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_side_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 2, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             "PROPOSED SOLUTION", font_size=12, bold=True, color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "JalSakhi: Any Smartphone into a Water Lab",
             font_size=34, bold=True, color=DEEP_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "A forensic water analysis platform with two complementary sensing modalities + community intelligence",
             font_size=16, color=GRAY_TEXT)

# Mode 1 Box
m1_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.0)
)
m1_box.fill.solid()
m1_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
m1_box.line.color.rgb = WATER_BLUE
m1_box.line.width = Pt(1.5)

add_text_box(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(0.5),
             "Mode 1: Electrochemical Fingerprinting",
             font_size=18, bold=True, color=WATER_BLUE)

add_text_box(slide, Inches(1.1), Inches(2.8), Inches(5), Inches(0.5),
             "PRECISION MODE", font_size=11, bold=True, color=TEAL)

add_bullet_list(slide, Inches(1.1), Inches(3.2), Inches(5), Inches(2.8), [
    "Pocket potentiostat dongle (INR 1,200) via BLE",
    "Disposable screen-printed electrodes (INR 25 each)",
    "Voltammetric scans: CV, DPV, SWV techniques",
    "On-device 1D-CNN classifies contaminants",
    "Results in under 60 seconds, fully offline",
    "Detects: Ammonia, Lead, Arsenic, Nitrate, Iron, Fluoride",
], font_size=14)

# Mode 2 Box
m2_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.5), Inches(4.0)
)
m2_box.fill.solid()
m2_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF0)
m2_box.line.color.rgb = RGBColor(0x2E, 0xA0, 0x6A)
m2_box.line.width = Pt(1.5)

add_text_box(slide, Inches(7.1), Inches(2.3), Inches(5), Inches(0.5),
             "Mode 2: Colorimetric Strip Analysis",
             font_size=18, bold=True, color=RGBColor(0x2E, 0xA0, 0x6A))

add_text_box(slide, Inches(7.1), Inches(2.8), Inches(5), Inches(0.5),
             "RAPID SCREENING MODE", font_size=11, bold=True, color=RGBColor(0x2E, 0xA0, 0x6A))

add_bullet_list(slide, Inches(7.1), Inches(3.2), Inches(5), Inches(2.8), [
    "Commercial test strips (INR 4-10 each)",
    "Smartphone camera captures color changes",
    "ArUco calibration card corrects for lighting",
    "CNN model quantifies concentrations from color",
    "Results in 30 seconds",
    "ZERO additional hardware needed",
], font_size=14)

# Platform layer at bottom
plat_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.6)
)
plat_box.fill.solid()
plat_box.fill.fore_color.rgb = DEEP_BLUE
plat_box.line.fill.background()

add_text_box(slide, Inches(1.1), Inches(6.42), Inches(11), Inches(0.5),
             "Platform Layer: Community Contamination Intelligence  |  Heatmaps  |  Temporal Forecasting  |  Municipal Dashboard",
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 3: SOLUTION ARCHITECTURE
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_side_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 3, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             "PROPOSED SOLUTION", font_size=12, bold=True, color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "System Architecture", font_size=34, bold=True, color=DEEP_BLUE)

# Three-layer architecture
layers = [
    {
        "title": "SENSING LAYER",
        "color": WATER_BLUE,
        "bg": RGBColor(0xE8, 0xF4, 0xFD),
        "items": [
            "Screen-Printed Electrode (3-electrode: WE/RE/CE)",
            "Potentiostat Dongle (ESP32 + Op-Amp circuit)",
            "CV / DPV / SWV waveform generation",
            "Paper Test Strips + Phone Camera + Calibration Card",
        ]
    },
    {
        "title": "EDGE AI LAYER (Smartphone)",
        "color": RGBColor(0x2E, 0xA0, 0x6A),
        "bg": RGBColor(0xE8, 0xF8, 0xF0),
        "items": [
            "Signal Processing: Savitzky-Golay smoothing, ALS baseline correction, peak detection",
            "ML Classification: 1D-CNN on voltammograms (TFLite, <200KB, <50ms)",
            "Colorimetric Analysis: Color correction + CNN regression",
            "Output: Contaminant ID + Concentration + Safety Rating + Treatment Advisory",
        ]
    },
    {
        "title": "CLOUD INTELLIGENCE LAYER",
        "color": RGBColor(0x9C, 0x27, 0xB0),
        "bg": RGBColor(0xF3, 0xE5, 0xF5),
        "items": [
            "Spatial Interpolation: Ordinary Kriging for contamination heatmaps",
            "Temporal Forecasting: LSTM on historical readings + seasonal patterns",
            "Anomaly Detection: Isolation Forest with source attribution",
            "Municipal Dashboard: District water safety scores + Jal Jeevan Mission integration",
        ]
    }
]

layer_y = Inches(1.5)
layer_h = Inches(1.7)
layer_gap = Inches(0.2)

for i, layer in enumerate(layers):
    y = layer_y + i * (layer_h + layer_gap)

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), layer_h
    )
    box.fill.solid()
    box.fill.fore_color.rgb = layer["bg"]
    box.line.color.rgb = layer["color"]
    box.line.width = Pt(1.5)

    # Layer title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(3.5), Inches(0.4)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = layer["color"]
    title_bar.line.fill.background()

    add_text_box(slide, Inches(1.0), y + Inches(0.02), Inches(3.3), Inches(0.35),
                 layer["title"], font_size=12, bold=True, color=WHITE)

    add_bullet_list(slide, Inches(1.1), y + Inches(0.45), Inches(10.8), layer_h - Inches(0.5),
                    layer["items"], font_size=13)

    # Arrow between layers
    if i < 2:
        arrow_y = y + layer_h
        add_text_box(slide, Inches(6), arrow_y, Inches(1), Inches(0.25),
                     "\u25BC", font_size=18, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 4: METHODOLOGY - ELECTROCHEMICAL SENSING
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_side_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 4, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             "METHODOLOGY", font_size=12, bold=True, color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Electrochemical Water Forensics", font_size=34, bold=True,
             color=DEEP_BLUE)

# How it works - left side
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
             "How It Works", font_size=20, bold=True, color=WATER_BLUE)

add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.0), [
    "Dip disposable electrode into water sample",
    "Potentiostat applies controlled voltage sweep",
    "Measures picoamp-to-milliamp current response",
    "Each contaminant produces unique electrochemical signature",
    "On-device CNN identifies & quantifies in 60 seconds",
    "Temperature-compensated (TMP117, \u00b10.1\u00b0C)",
    "Auto-range TIA: 10 nA to 10 mA dynamic range",
], font_size=14)

# Detection table - right side
add_text_box(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.4),
             "Detection Capabilities", font_size=20, bold=True, color=WATER_BLUE)

table_data = [
    ["Contaminant", "Detection Limit", "Method"],
    ["Ammonia", "0.05 mg/L", "SWV (Prussian Blue SPE)"],
    ["Lead", "1 ppb", "DPASV (Bismuth-film)"],
    ["Arsenic", "5 ppb", "DPASV (Gold nanoparticle)"],
    ["Nitrate", "0.5 mg/L", "CV (Copper-modified)"],
    ["Iron", "0.05 mg/L", "DPV (bare carbon)"],
    ["Fluoride", "0.1 mg/L", "Potentiometry (ISE)"],
    ["Free Chlorine", "0.1 mg/L", "Amperometry"],
]

make_table(slide, Inches(6.8), Inches(2.0), Inches(5.5), table_data,
           font_size=12)

# Key specs box at bottom
spec_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.4)
)
spec_box.fill.solid()
spec_box.fill.fore_color.rgb = LIGHT_TEAL
spec_box.line.color.rgb = TEAL
spec_box.line.width = Pt(1)

add_text_box(slide, Inches(1.1), Inches(5.55), Inches(4), Inches(0.35),
             "Prototype Hardware (INR 4,000 Budget)", font_size=14, bold=True,
             color=DEEP_BLUE)

proto_items = [
    "ESP32 dev board (BLE + WiFi + DAC + ADC) - INR 500",
    "LM358 op-amps on breadboard - INR 30",
    "DIY electrodes: pencil graphite (WE/CE) + silver wire AgCl (RE)",
    "Published technique: pencil graphite as carbon electrode (Electrochimica Acta)",
]
add_bullet_list(slide, Inches(1.1), Inches(5.9), Inches(10.8), Inches(0.9),
                proto_items, font_size=12)


# ══════════════════════════════════════════════════════════
# SLIDE 5: METHODOLOGY - AI PIPELINE
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_side_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 5, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             "METHODOLOGY", font_size=12, bold=True, color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "AI & Machine Learning Pipeline", font_size=34, bold=True,
             color=DEEP_BLUE)

# Signal Processing - Left
sp_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5)
)
sp_box.fill.solid()
sp_box.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
sp_box.line.color.rgb = ORANGE
sp_box.line.width = Pt(1)

add_text_box(slide, Inches(1.1), Inches(1.55), Inches(5), Inches(0.4),
             "On-Device Signal Processing", font_size=18, bold=True,
             color=ORANGE)

add_bullet_list(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(1.8), [
    "Savitzky-Golay smoothing (order 3, window 15)",
    "ALS baseline correction (removes capacitive bg)",
    "Derivative-based peak detection (5x noise threshold)",
    "Feature extraction: Ep, Ip, half-width, area",
    "Quality: GOOD / MARGINAL / REJECT",
], font_size=13)

# ML Model - Right
ml_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.5), Inches(2.5)
)
ml_box.fill.solid()
ml_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
ml_box.line.color.rgb = WATER_BLUE
ml_box.line.width = Pt(1)

add_text_box(slide, Inches(7.1), Inches(1.55), Inches(5), Inches(0.4),
             "1D-CNN Classification Model", font_size=18, bold=True,
             color=WATER_BLUE)

add_bullet_list(slide, Inches(7.1), Inches(2.0), Inches(5), Inches(1.8), [
    "Input: 1000-pt voltammogram + metadata (TDS, pH, temp)",
    "Conv1D(32,k=7) -> Conv1D(64,k=5) -> Conv1D(128,k=3)",
    "Dual heads: Detection (sigmoid) + Concentration (ReLU)",
    "INT8 quantized TFLite: <200 KB, <50 ms inference",
    "Fully offline, no internet required",
], font_size=13)

# Colorimetric - Left bottom
cv_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.0)
)
cv_box.fill.solid()
cv_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF0)
cv_box.line.color.rgb = RGBColor(0x2E, 0xA0, 0x6A)
cv_box.line.width = Pt(1)

add_text_box(slide, Inches(1.1), Inches(4.35), Inches(5), Inches(0.4),
             "Colorimetric AI Pipeline", font_size=18, bold=True,
             color=RGBColor(0x2E, 0xA0, 0x6A))

add_bullet_list(slide, Inches(1.1), Inches(4.8), Inches(5), Inches(1.3), [
    "ArUco marker detection + perspective correction",
    "10-patch calibration card for color/lighting correction",
    "CNN regression on LAB color space",
    "Cross-phone accuracy: <5% deltaE after correction",
], font_size=13)

# Confidence & Safety - Right bottom
conf_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.3), Inches(5.5), Inches(2.0)
)
conf_box.fill.solid()
conf_box.fill.fore_color.rgb = RGBColor(0xF3, 0xE5, 0xF5)
conf_box.line.color.rgb = RGBColor(0x9C, 0x27, 0xB0)
conf_box.line.width = Pt(1)

add_text_box(slide, Inches(7.1), Inches(4.35), Inches(5), Inches(0.4),
             "Confidence & Safety Scoring", font_size=18, bold=True,
             color=RGBColor(0x9C, 0x27, 0xB0))

add_bullet_list(slide, Inches(7.1), Inches(4.8), Inches(5), Inches(1.3), [
    "Model probability + SNR + scan quality flag",
    "Output: HIGH / MEDIUM / LOW / RETEST",
    "Autoencoder anomaly detector flags interference",
    "Prevents false safety signals",
], font_size=13)

# Key stat at bottom
add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
             "Synthetic training data: Physics-based voltammogram generator (Gaussian peaks + Randles-Sevcik kinetics) | Domain adaptation for water matrix types",
             font_size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 6: NOVELTY - WHAT MAKES THIS DIFFERENT
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_side_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 6, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             "NOVELTY", font_size=12, bold=True, color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Why This Is Fundamentally Different",
             font_size=34, bold=True, color=DEEP_BLUE)

# Comparison table
comp_data = [
    ["Aspect", "Existing Solutions", "JalSakhi"],
    ["Cost", "Fixed IoT sensors: $200-$2,000", "Portable dongle: $14 + $0.30/test"],
    ["Parameters", "3-4 parameters continuously", "7+ contaminants incl. heavy metals at ppb"],
    ["Data", "Siloed per household", "Crowdsourced contamination intelligence maps"],
    ["Calibration", "Sensors degrade, drift, need recalibration", "Disposable electrodes = fresh every test"],
    ["Hardware", "Dedicated hardware + internet", "Any smartphone, fully offline"],
    ["Approach", "Passive monitoring at one point", "Active forensic testing of any source"],
]

make_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), comp_data, font_size=13)

# Key technical insight
insight_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.0)
)
insight_box.fill.solid()
insight_box.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
insight_box.line.color.rgb = ORANGE
insight_box.line.width = Pt(1.5)

add_text_box(slide, Inches(1.1), Inches(4.85), Inches(5), Inches(0.4),
             "Key Technical Insight", font_size=20, bold=True, color=ORANGE)

novelties = [
    "Disposable sensing elements eliminate drift & calibration entirely - fresh calibration every test",
    "Smartphone is the instrument - leverages $100+ computing power users already own",
    "Human-in-the-loop: women actively test suspicious sources, not passive monitoring",
    "Multi-point data: one device tests hundreds of sources across a community",
    "Community intelligence layer: individual tests aggregate into district-level contamination maps",
    "Treatment advisory: AI prescribes minimum effective treatment, not just detection",
]
add_bullet_list(slide, Inches(1.1), Inches(5.3), Inches(10.8), Inches(1.4),
                novelties, font_size=13)


# ══════════════════════════════════════════════════════════
# SLIDE 7: ALL 5 CATEGORIES COVERED
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_side_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 7, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             "NOVELTY", font_size=12, bold=True, color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "One Platform, All Five Categories",
             font_size=34, bold=True, color=DEEP_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
             "Most teams pick one category. JalSakhi covers all five with a unified platform.",
             font_size=16, color=GRAY_TEXT)

categories = [
    ("1. Ammonia Mitigation", "SWV on Prussian Blue SPE | 0.05 mg/L detection (10x WHO limit) | AI-generated dosage protocol", WATER_BLUE),
    ("2. Water Neutrality", "Community water budgeting | Contaminated water avoided via informed source selection | Conservation tracking", RGBColor(0x2E, 0xA0, 0x6A)),
    ("3. Smart Distribution", "Crowdsourced data maps pipeline contamination | Chlorine residual drop = intrusion | Safest source routing", RGBColor(0x9C, 0x27, 0xB0)),
    ("4. Affordable Sensing", "INR 1,200 dongle + INR 25/test | 7+ contaminants incl. heavy metals | Colorimetric mode = ZERO hardware", ORANGE),
    ("5. AI Water Management", "Edge CNN (offline) + Cloud Kriging/LSTM/Anomaly Detection | Municipal decision support dashboard", RED_ACCENT),
]

cat_y = Inches(1.8)
cat_h = Inches(0.9)
cat_gap = Inches(0.15)

for i, (title, desc, color) in enumerate(categories):
    y = cat_y + i * (cat_h + cat_gap)

    # Category box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), cat_h
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(1.5)

    # Color bar on left
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.15), cat_h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text_box(slide, Inches(1.2), y + Inches(0.08), Inches(10.8), Inches(0.35),
                 title, font_size=16, bold=True, color=color)
    add_text_box(slide, Inches(1.2), y + Inches(0.42), Inches(10.8), Inches(0.4),
                 desc, font_size=13, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════
# SLIDE 8: RESOURCES REQUIRED & BUDGET
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_side_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 8, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             "RESOURCES REQUIRED", font_size=12, bold=True, color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Budget & Resources", font_size=34, bold=True, color=DEEP_BLUE)

# Competition Prototype Budget - Left
proto_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.0)
)
proto_box.fill.solid()
proto_box.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
proto_box.line.color.rgb = ORANGE
proto_box.line.width = Pt(1.5)

add_text_box(slide, Inches(1.1), Inches(1.55), Inches(5), Inches(0.4),
             "Competition Prototype: INR 4,000", font_size=18, bold=True,
             color=ORANGE)

proto_table = [
    ["Item", "Cost (INR)"],
    ["ESP32 dev board (BLE+WiFi+DAC+ADC)", "500"],
    ["Water test strips (16-in-1, 100 pack)", "400"],
    ["LM358 op-amps (x2)", "30"],
    ["Breadboard + jumper wires", "250"],
    ["Resistor assortment", "100"],
    ["Silver wire (reference electrode)", "200"],
    ["Ammonia solution (demo)", "100"],
    ["KCl electrolyte", "50"],
    ["Calibration card + posters", "210"],
    ["Total", "1,840"],
    ["Buffer", "2,160"],
]

make_table(slide, Inches(1.1), Inches(2.0), Inches(5), proto_table, font_size=11)

# Production Scale - Right
prod_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.5), Inches(4.0)
)
prod_box.fill.solid()
prod_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
prod_box.line.color.rgb = WATER_BLUE
prod_box.line.width = Pt(1.5)

add_text_box(slide, Inches(7.1), Inches(1.55), Inches(5), Inches(0.4),
             "Production Scale Economics", font_size=18, bold=True,
             color=WATER_BLUE)

scale_table = [
    ["Scale", "SPE Cost", "Kit Cost", "Cost/Test"],
    ["Prototype (10 kits)", "$0.80", "$25", "$1.50"],
    ["Pilot (1,000 kits)", "$0.40", "$15", "$0.60"],
    ["Scale (100K kits)", "$0.20", "$10", "$0.30"],
]
make_table(slide, Inches(7.1), Inches(2.0), Inches(5), scale_table, font_size=12)

add_text_box(slide, Inches(7.1), Inches(3.4), Inches(5), Inches(0.4),
             "Production BOM (AD5940 + STM32)", font_size=14, bold=True,
             color=WATER_BLUE)

prod_bom = [
    ["Component", "Cost"],
    ["AD5940 AFE", "$4.50"],
    ["STM32L432 MCU", "$2.50"],
    ["TMP117 temp sensor", "$1.20"],
    ["Pogo pins (3x Mill-Max)", "$1.20"],
    ["LDOs, ferrites, caps, ESD", "$0.85"],
    ["4-layer PCB + shield", "$1.50"],
    ["Assembly + enclosure", "$1.90"],
    ["Total Dongle", "$14.29 (INR ~1,200)"],
]
make_table(slide, Inches(7.1), Inches(3.8), Inches(5), prod_bom, font_size=11)

# Key stat
key_stat = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0)
)
key_stat.fill.solid()
key_stat.fill.fore_color.rgb = DEEP_BLUE
key_stat.line.fill.background()

add_text_box(slide, Inches(1.1), Inches(5.85), Inches(11), Inches(0.4),
             "INR 25 per test  vs  INR 500-2,000 lab test  =  80-95% cost reduction",
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.1), Inches(6.3), Inches(11), Inches(0.35),
             "Every part number verified against Mouser/DigiKey. Every component commercially available.",
             font_size=13, color=TEAL, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 9: RESULTS ACHIEVED
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_side_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 9, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             "RESULTS ACHIEVED", font_size=12, bold=True, color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "What We Have Built & Validated", font_size=34, bold=True,
             color=DEEP_BLUE)

# Result 1: Technical Validation
r1_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.3)
)
r1_box.fill.solid()
r1_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
r1_box.line.color.rgb = WATER_BLUE
r1_box.line.width = Pt(1)

add_text_box(slide, Inches(1.1), Inches(1.55), Inches(5), Inches(0.4),
             "Technical Feasibility - Validated", font_size=18, bold=True,
             color=WATER_BLUE)

add_bullet_list(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(1.6), [
    "All components commercially available, sourced",
    "Published DPV on SPEs: >95% correlation with ICP-MS",
    "Pencil graphite electrodes: published technique",
    "ESP32 + op-amp circuit functional for CV/DPV",
    "Prototype built within INR 2,000 budget",
], font_size=13)

# Result 2: AI Validation
r2_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.5), Inches(2.3)
)
r2_box.fill.solid()
r2_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF0)
r2_box.line.color.rgb = RGBColor(0x2E, 0xA0, 0x6A)
r2_box.line.width = Pt(1)

add_text_box(slide, Inches(7.1), Inches(1.55), Inches(5), Inches(0.4),
             "AI Pipeline - Designed & Specified", font_size=18, bold=True,
             color=RGBColor(0x2E, 0xA0, 0x6A))

add_bullet_list(slide, Inches(7.1), Inches(2.0), Inches(5), Inches(1.6), [
    "1D-CNN architecture: domain-adapted, multi-task",
    "Synthetic data generation pipeline specified",
    "Signal processing: peer-reviewed algorithms",
    "TFLite quantization: <200KB model, <50ms",
    "Confidence scoring with interference detection",
], font_size=13)

# Result 3: Literature Backing
r3_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.1), Inches(5.5), Inches(2.5)
)
r3_box.fill.solid()
r3_box.fill.fore_color.rgb = RGBColor(0xF3, 0xE5, 0xF5)
r3_box.line.color.rgb = RGBColor(0x9C, 0x27, 0xB0)
r3_box.line.width = Pt(1)

add_text_box(slide, Inches(1.1), Inches(4.15), Inches(5), Inches(0.4),
             "Backed by Published Research", font_size=18, bold=True,
             color=RGBColor(0x9C, 0x27, 0xB0))

add_bullet_list(slide, Inches(1.1), Inches(4.6), Inches(5), Inches(1.8), [
    "26 peer-reviewed references supporting every claim",
    "Smartphone potentiostats: Ainla 2018, Nemiroski 2014",
    "Electrochemical water analysis: Cui 2015, Li 2016",
    "ML for electrochemistry: Kammarchedu 2022",
    "Community monitoring: Kohlitz 2020, Zheng 2019",
    "India context: Jal Jeevan Mission, CGWB 2018",
], font_size=13)

# Result 4: Key numbers
r4_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.1), Inches(5.5), Inches(2.5)
)
r4_box.fill.solid()
r4_box.fill.fore_color.rgb = RGBColor(0xFE, 0xF9, 0xE7)
r4_box.line.color.rgb = ORANGE
r4_box.line.width = Pt(1)

add_text_box(slide, Inches(7.1), Inches(4.15), Inches(5), Inches(0.4),
             "Key Performance Targets", font_size=18, bold=True,
             color=ORANGE)

kpi_data = [
    ["Metric", "Target"],
    ["Ammonia detection limit", "0.05 mg/L (10x WHO)"],
    ["Lead detection limit", "1 ppb"],
    ["Classification time", "<60 seconds"],
    ["Model size", "<200 KB (TFLite)"],
    ["Cost per test", "INR 25"],
    ["Lab correlation target", ">90%"],
    ["Dongle BOM", "$14.29 (INR 1,200)"],
]
make_table(slide, Inches(7.1), Inches(4.6), Inches(5), kpi_data, font_size=12)


# ══════════════════════════════════════════════════════════
# SLIDE 10: DEPLOYMENT MODEL - WOMEN AS INFRASTRUCTURE
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_side_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 10, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             "CASE STUDY & DEPLOYMENT", font_size=12, bold=True, color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Women as Water Quality Infrastructure",
             font_size=34, bold=True, color=DEEP_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
             "Leveraging India's 12 million women Self-Help Groups (NRLM) - already organized, trained, smartphone-equipped",
             font_size=15, color=GRAY_TEXT)

# Deployment levels
levels = [
    ("SHG Level", "12M SHGs, 140M members, 99% block coverage", "2-3 Jal Sakhis per SHG | Weekly testing of community sources | INR 10-20/validated test via UPI", WATER_BLUE),
    ("Village Level", "Gram Panchayat Water Quality Register", "All test data in village water register | Paani Samiti integration | Evidence-based infrastructure decisions", RGBColor(0x2E, 0xA0, 0x6A)),
    ("District Level", "Block/District Intelligence Dashboard", "Contamination heatmaps | Integration with JJM IMIS | Real-time dashboards for District Water Officers", RGBColor(0x9C, 0x27, 0xB0)),
]

level_y = Inches(1.8)
for i, (title, subtitle, desc, color) in enumerate(levels):
    y = level_y + i * Inches(1.5)

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(1.3)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(1.5)

    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.25), Inches(0.7), Inches(0.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text_box(slide, Inches(1.0), y + Inches(0.3), Inches(0.7), Inches(0.6),
                 str(i + 1), font_size=24, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2.0), y + Inches(0.1), Inches(10), Inches(0.35),
                 title, font_size=16, bold=True, color=color)
    add_text_box(slide, Inches(2.0), y + Inches(0.4), Inches(10), Inches(0.3),
                 subtitle, font_size=12, bold=True, color=GRAY_TEXT)
    add_text_box(slide, Inches(2.0), y + Inches(0.7), Inches(10), Inches(0.5),
                 desc, font_size=13, color=DARK_TEXT)

# Deployment phases at bottom
phase_y = Inches(6.3)
phases = [
    ("Phase 1: PoC", "5 SHGs, 1 GP\n3 months, INR 2L", WATER_BLUE),
    ("Phase 2: Pilot", "100 SHGs, 3 blocks\n6 months, INR 10L", RGBColor(0x2E, 0xA0, 0x6A)),
    ("Phase 3: Scale", "5,000 SHGs, 1 state\n12 months, INR 1 Cr", RGBColor(0x9C, 0x27, 0xB0)),
]

phase_w = Inches(3.4)
for i, (title, desc, color) in enumerate(phases):
    x = Inches(0.8) + i * (phase_w + Inches(0.3))
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, phase_y, phase_w, Inches(0.7)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    add_text_box(slide, x + Inches(0.1), phase_y + Inches(0.02), Inches(1.3), Inches(0.3),
                 title, font_size=11, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(1.4), phase_y + Inches(0.02), Inches(1.8), Inches(0.6),
                 desc, font_size=10, color=WHITE)


# ══════════════════════════════════════════════════════════
# SLIDE 11: IMPACT QUANTIFICATION
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_side_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 11, TOTAL_SLIDES)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(4), Inches(0.4),
             "KEY INFERENCES", font_size=12, bold=True, color=TEAL)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Impact Mapped to the Problem",
             font_size=34, bold=True, color=DEEP_BLUE)

# Impact cards
impacts = [
    {
        "icon": "HEALTH",
        "title": "200-500 Hospitalizations Prevented/Year",
        "desc": "15-20% sources with ammonia > WHO limit detected\n~5,000 families alerted to unsafe water",
        "color": RED_ACCENT,
        "bg": RGBColor(0xFD, 0xE8, 0xE8),
    },
    {
        "icon": "TIME",
        "title": "450,000 Hours Saved Annually",
        "desc": "Contamination maps show safe sources\n15 min avg reduction in collection time per trip",
        "color": WATER_BLUE,
        "bg": RGBColor(0xE8, 0xF4, 0xFD),
    },
    {
        "icon": "INCOME",
        "title": "INR 72 Lakh/Year into Women's Hands",
        "desc": "1,000 Jal Sakhis x 10 tests/week x INR 15\nSkills: STEM, data collection, community health",
        "color": RGBColor(0x2E, 0xA0, 0x6A),
        "bg": RGBColor(0xE8, 0xF8, 0xF0),
    },
    {
        "icon": "GOVERNANCE",
        "title": "52x Monitoring Frequency Increase",
        "desc": "Weekly testing vs annual lab surveys\nReal-time ground truth for Jal Jeevan Mission",
        "color": RGBColor(0x9C, 0x27, 0xB0),
        "bg": RGBColor(0xF3, 0xE5, 0xF5),
    },
]

card_w = Inches(5.5)
card_h = Inches(1.2)
for i, imp in enumerate(impacts):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * (card_w + Inches(0.5))
    y = Inches(1.5) + row * (card_h + Inches(0.3))

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h
    )
    box.fill.solid()
    box.fill.fore_color.rgb = imp["bg"]
    box.line.color.rgb = imp["color"]
    box.line.width = Pt(1.5)

    # Icon label
    add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(1), Inches(0.3),
                 imp["icon"], font_size=10, bold=True, color=imp["color"])
    add_text_box(slide, x + Inches(0.15), y + Inches(0.3), card_w - Inches(0.3), Inches(0.35),
                 imp["title"], font_size=15, bold=True, color=imp["color"])
    add_text_box(slide, x + Inches(0.15), y + Inches(0.65), card_w - Inches(0.3), Inches(0.5),
                 imp["desc"], font_size=12, color=DARK_TEXT)

# Government alignment at bottom
gov_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.3)
)
gov_box.fill.solid()
gov_box.fill.fore_color.rgb = LIGHT_GRAY
gov_box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
gov_box.line.width = Pt(1)

add_text_box(slide, Inches(1.1), Inches(4.55), Inches(11), Inches(0.4),
             "Alignment with Government Programs & SDGs", font_size=18, bold=True,
             color=DEEP_BLUE)

gov_items = [
    "Jal Jeevan Mission: Continuous quality monitoring after tap installation (filling the gap JJM currently has)",
    "NRLM / DAY-NRLM: Leverages existing SHG infrastructure + creates new livelihood for women",
    "Swachh Bharat Mission - Gramin: Water quality data supports ODF Plus monitoring (Phase 2)",
    "Atal Bhujal Yojana: Community participation in groundwater monitoring across 7 states",
    "SDG Alignment: SDG 3 (Health), SDG 5 (Gender), SDG 6 (Water), SDG 8 (Work), SDG 9 (Innovation)",
]
add_bullet_list(slide, Inches(1.1), Inches(5.0), Inches(10.8), Inches(1.7),
                gov_items, font_size=13)

# L&T angle
add_text_box(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.3),
             "For L&T: This platform tells you WHERE to build treatment plants. Contamination intelligence at scale = demand signal for infrastructure.",
             font_size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 12: THE ASK
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg_rect(slide, DEEP_BLUE)
add_accent_bar(slide, y=Inches(0), height=Inches(0.08), color=TEAL)

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.6),
             "The Ask", font_size=40, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

# Pilot details box
pilot_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.8), Inches(9), Inches(3.2)
)
pilot_box.fill.solid()
pilot_box.fill.fore_color.rgb = RGBColor(0x15, 0x3A, 0x5E)
pilot_box.line.color.rgb = TEAL
pilot_box.line.width = Pt(2)

pilot_stats = [
    ("100 SHGs", "across 3 ammonia-affected districts"),
    ("6 Months", "proof-of-value pilot"),
    ("INR 5 Lakh", "total investment"),
]

for i, (num, desc) in enumerate(pilot_stats):
    y = Inches(2.0) + i * Inches(0.9)
    add_text_box(slide, Inches(3), y, Inches(3), Inches(0.5),
                 num, font_size=32, bold=True, color=TEAL,
                 align=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(6.2), y + Inches(0.05), Inches(4), Inches(0.5),
                 desc, font_size=18, color=WHITE)

# Deliverable
add_text_box(slide, Inches(2), Inches(5.2), Inches(9), Inches(0.5),
             "Deliverable", font_size=14, bold=True, color=TEAL,
             align=PP_ALIGN.CENTER)

del_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.6), Inches(8), Inches(0.6)
)
del_box.fill.solid()
del_box.fill.fore_color.rgb = TEAL
del_box.line.fill.background()

add_text_box(slide, Inches(2.5), Inches(5.62), Inches(8), Inches(0.55),
             "Validated contamination intelligence dashboard\nwith >90% correlation to NABL lab results",
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom tagline
add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             '"Women become the sensing infrastructure of India\'s water system"',
             font_size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Credits
add_text_box(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.3),
             "Every claim backed by 26 peer-reviewed references. All components commercially available. Prototype demonstrable.",
             font_size=11, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

add_accent_bar(slide, y=Inches(7.42), height=Inches(0.08), color=TEAL)


# ── SAVE ──────────────────────────────────────────────────
filename = f"{TEAM_NAME}_{INSTITUTE_NAME}.pptx".replace(" ", "_")
output_path = f"C:/Users/Ujjwal/JalSakhi/presentation/{filename}"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: Title + {TOTAL_SLIDES} content slides")
print(f"File: {filename}")
